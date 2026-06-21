from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
ROOT = OUT_DIR.parents[1]
LOGO = ROOT / "public" / "brand" / "easyphone-logo.png"
PPTX = OUT_DIR / "easyphone-ai-roadshow.pptx"

WIDE = (13.333, 7.5)

NAVY = RGBColor(16, 36, 55)
BLUE = RGBColor(47, 128, 237)
MINT = RGBColor(190, 232, 211)
CREAM = RGBColor(255, 247, 200)
RED = RGBColor(255, 107, 95)
INK = RGBColor(35, 45, 56)
MUTED = RGBColor(105, 116, 128)
PAPER = RGBColor(249, 251, 247)
LINE = RGBColor(223, 229, 235)
WHITE = RGBColor(255, 255, 255)


def add_textbox(slide, x, y, w, h, text, size=24, color=INK, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.72, 0.52, 8.6, 0.6, title, size=30, color=NAVY, bold=True)
    if subtitle:
        add_textbox(slide, 0.74, 1.12, 8.7, 0.38, subtitle, size=12, color=MUTED)


def add_footer(slide, idx):
    add_textbox(slide, 0.72, 7.05, 4, 0.24, "EasyPhone AI / 爸妈别急", size=9, color=MUTED)
    add_textbox(slide, 12.2, 7.05, 0.4, 0.24, f"{idx:02d}", size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_round_rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_bullet_list(slide, x, y, w, h, items, size=18, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
    return box


def add_chip(slide, x, y, text, fill, color=INK, w=None):
    width = w if w is not None else max(1.1, 0.25 + len(text) * 0.16)
    add_round_rect(slide, x, y, width, 0.34, fill, radius=True)
    add_textbox(slide, x + 0.1, y + 0.06, width - 0.2, 0.18, text, size=9, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_big_number(slide, x, y, number, title, body, color):
    add_textbox(slide, x, y, 0.7, 0.45, number, size=24, color=color, bold=True)
    add_textbox(slide, x + 0.7, y + 0.03, 2.9, 0.36, title, size=16, color=NAVY, bold=True)
    add_textbox(slide, x + 0.7, y + 0.48, 3.1, 0.7, body, size=11, color=MUTED)


def set_bg(slide, color=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(WIDE[0])
    prs.slide_height = Inches(WIDE[1])
    blank = prs.slide_layouts[6]

    # 1 Cover
    slide = prs.slides.add_slide(blank)
    set_bg(slide, PAPER)
    add_round_rect(slide, 0, 0, 13.333, 7.5, PAPER, radius=False)
    add_round_rect(slide, 0, 0, 4.15, 7.5, MINT, radius=False)
    add_round_rect(slide, 0.76, 0.72, 1.25, 1.25, WHITE, radius=True)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.88), Inches(0.84), Inches(1.0), Inches(1.0))
    add_textbox(slide, 4.85, 1.12, 6.8, 0.5, "EasyPhone AI", size=22, color=BLUE, bold=True)
    add_textbox(slide, 4.8, 1.75, 6.9, 1.0, "爸妈别急", size=46, color=NAVY, bold=True)
    add_textbox(slide, 4.85, 2.72, 6.7, 0.48, "给低识字中老年人的 AI 语音手机教练", size=18, color=INK)
    add_textbox(slide, 4.85, 3.55, 6.7, 0.72, "普通问题一步一步教，危险问题先停下来。", size=27, color=RED, bold=True)
    add_chip(slide, 4.85, 4.62, "AI Agent", CREAM, NAVY, w=1.2)
    add_chip(slide, 6.18, 4.62, "Elderly Digital Safety", WHITE, NAVY, w=2.2)
    add_chip(slide, 8.55, 4.62, "Human-centered AI", WHITE, NAVY, w=2.0)
    add_textbox(slide, 4.85, 6.55, 5.8, 0.25, "Demo: https://easy-phone-ai.vercel.app", size=10, color=MUTED)
    add_textbox(slide, 4.85, 6.85, 5.8, 0.25, "GitHub: github.com/qrx-joe/EasyPhone_AI", size=10, color=MUTED)

    # 2 Problem
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "问题：智能手机对老人是高压系统", "他们不是不想学，而是看不懂、说不清、也不敢点。")
    add_round_rect(slide, 0.72, 1.72, 3.75, 4.55, WHITE, LINE)
    add_round_rect(slide, 4.8, 1.72, 3.75, 4.55, WHITE, LINE)
    add_round_rect(slide, 8.88, 1.72, 3.75, 4.55, WHITE, LINE)
    add_textbox(slide, 1.02, 2.08, 3.0, 0.36, "老人端", size=20, color=NAVY, bold=True)
    add_bullet_list(slide, 1.02, 2.75, 3.0, 2.7, ["不会打字和搜索", "看不懂弹窗和按钮", "害怕误删、扣费、被骗"], size=15)
    add_textbox(slide, 5.1, 2.08, 3.0, 0.36, "高风险场景", size=20, color=NAVY, bold=True)
    add_bullet_list(slide, 5.1, 2.75, 3.05, 2.7, ["验证码、转账", "陌生链接、二维码", "医保短信、屏幕共享"], size=15)
    add_textbox(slide, 9.18, 2.08, 3.0, 0.36, "子女端", size=20, color=NAVY, bold=True)
    add_bullet_list(slide, 9.18, 2.75, 3.05, 2.7, ["电话指导效率低", "父母描述不清", "担心诈骗发生在身边之外"], size=15)
    add_textbox(slide, 1.0, 6.55, 11.4, 0.36, "真正的缺口不是“更多教程”，而是一个能慢慢教、也敢及时停下来的安全教练。", size=18, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 2)

    # 3 Solution
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "方案：一个有安全边界的手机教练", "不是替老人操作手机，而是帮助老人安全地理解下一步。")
    steps = [
        ("1", "语音提问", "老人直接说“微信没有声音了”。"),
        ("2", "风险判断", "规则兜底 + AI 增强，先判断能不能教。"),
        ("3", "低风险分步教", "每次只给一步，大字 + 语音播报。"),
        ("4", "高风险立刻停", "不进入教程，生成家人求助卡。"),
    ]
    x_positions = [0.82, 3.95, 7.08, 10.21]
    for x, (num, title, body) in zip(x_positions, steps):
        add_round_rect(slide, x, 2.05, 2.55, 3.7, WHITE, LINE)
        add_textbox(slide, x + 0.25, 2.32, 0.55, 0.42, num, size=24, color=BLUE, bold=True)
        add_textbox(slide, x + 0.25, 3.05, 2.05, 0.36, title, size=18, color=NAVY, bold=True)
        add_textbox(slide, x + 0.25, 3.75, 2.05, 1.05, body, size=13, color=MUTED)
    add_round_rect(slide, 1.05, 6.22, 11.25, 0.64, NAVY, radius=True)
    add_textbox(slide, 1.35, 6.39, 10.6, 0.24, "核心闭环：老人说问题 → 系统先判风险 → 能教才教，危险就停", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 3)

    # 4 Differentiation
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "关键差异：AI 知道什么时候不该回答", "反诈场景里，乱帮忙比不帮忙更危险。")
    add_textbox(slide, 0.85, 1.86, 5.2, 0.8, "Most AI assistants try to answer every question.", size=25, color=MUTED, bold=True)
    add_textbox(slide, 0.85, 2.95, 5.4, 1.05, "EasyPhone AI knows when not to answer.", size=34, color=RED, bold=True)
    add_round_rect(slide, 7.05, 1.58, 5.35, 4.8, WHITE, LINE)
    add_bullet_list(slide, 7.45, 2.0, 4.6, 3.7, [
        "先判断风险，再决定是否指导",
        "规则兜底优先，AI 只做增强",
        "多关键词命中取最高风险等级",
        "高风险不进入普通教程",
        "求助卡不诱导老人泄露敏感信息",
    ], size=16)
    add_round_rect(slide, 0.85, 5.38, 5.42, 0.72, CREAM, radius=True)
    add_textbox(slide, 1.15, 5.58, 4.8, 0.28, "不是“更会答”，而是“更知道边界”。", size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)

    # 5 Flow
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "产品流程：低风险教，高风险停", "同一个入口，根据风险走两条完全不同的路径。")
    add_round_rect(slide, 0.9, 1.75, 2.1, 0.7, WHITE, LINE)
    add_textbox(slide, 1.1, 1.96, 1.7, 0.22, "用户输入", size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_round_rect(slide, 3.45, 1.75, 2.35, 0.7, WHITE, LINE)
    add_textbox(slide, 3.65, 1.96, 1.95, 0.22, "风险分类", size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_round_rect(slide, 6.35, 1.12, 2.6, 0.7, MINT, LINE)
    add_textbox(slide, 6.58, 1.33, 2.1, 0.22, "低风险", size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_round_rect(slide, 9.55, 1.12, 2.75, 0.7, WHITE, LINE)
    add_textbox(slide, 9.75, 1.33, 2.35, 0.22, "确认 + 分步指导", size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_round_rect(slide, 6.35, 3.52, 2.6, 0.7, RED, None)
    add_textbox(slide, 6.58, 3.73, 2.1, 0.22, "高风险", size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_round_rect(slide, 9.55, 3.52, 2.75, 0.7, WHITE, LINE)
    add_textbox(slide, 9.75, 3.73, 2.35, 0.22, "停止 + 求助卡", size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    # connector lines
    for x1, y1, x2, y2 in [(3.0, 2.1, 3.45, 2.1), (5.8, 2.1, 6.35, 1.47), (5.8, 2.1, 6.35, 3.87), (8.95, 1.47, 9.55, 1.47), (8.95, 3.87, 9.55, 3.87)]:
        line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        line.line.color.rgb = NAVY
        line.line.width = Pt(2)
    add_big_number(slide, 0.95, 5.25, "01", "每次只教一步", "降低阅读压力和误操作概率。", BLUE)
    add_big_number(slide, 4.75, 5.25, "02", "危险不继续教", "验证码、转账、屏幕共享立即停止。", RED)
    add_big_number(slide, 8.65, 5.25, "03", "家人能看懂", "把老人模糊表达整理成求助单。", BLUE)
    add_footer(slide, 5)

    # 6 Demo
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "Demo：3 个场景跑完整闭环", "路演只看一件事：系统何时教，何时停。")
    demos = [
        ("微信没有声音", "低风险", "确认问题 → 进入分步教程 → 语音播报", MINT),
        ("手机字体太小", "低风险", "大字短句 → 没看到可换一种说法", MINT),
        ("医保短信 + 验证码", "高风险", "停止指导 → 不点链接 → 生成家人求助卡", RGBColor(255, 220, 216)),
    ]
    for i, (title, risk, flow, color) in enumerate(demos):
        y = 1.68 + i * 1.48
        add_round_rect(slide, 0.92, y, 11.5, 1.08, WHITE, LINE)
        add_round_rect(slide, 1.18, y + 0.22, 1.35, 0.42, color, radius=True)
        add_textbox(slide, 1.32, y + 0.33, 1.05, 0.16, risk, size=9, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, 2.85, y + 0.2, 2.6, 0.32, title, size=19, color=NAVY, bold=True)
        add_textbox(slide, 5.65, y + 0.23, 5.9, 0.3, flow, size=15, color=MUTED)
    add_footer(slide, 6)

    # 7 Tech
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "技术与安全：规则兜底，AI 增强", "安全判断不能完全交给模型。")
    add_round_rect(slide, 0.82, 1.72, 5.55, 4.6, WHITE, LINE)
    add_textbox(slide, 1.15, 2.02, 4.8, 0.34, "技术栈", size=20, color=NAVY, bold=True)
    add_bullet_list(slide, 1.15, 2.65, 4.75, 2.6, [
        "Next.js 16 + React 19 + TypeScript",
        "Tailwind CSS 4",
        "Web Speech API 语音输入",
        "SpeechSynthesis 语音播报",
        "可选 DeepSeek AI 风险复检",
    ], size=14)
    add_round_rect(slide, 6.95, 1.72, 5.55, 4.6, WHITE, LINE)
    add_textbox(slide, 7.28, 2.02, 4.8, 0.34, "安全不变量", size=20, color=NAVY, bold=True)
    add_bullet_list(slide, 7.28, 2.65, 4.75, 2.6, [
        "高风险统一由路由层分流",
        "多关键词命中取最高风险 MAX(level)",
        "高风险不进入普通教程",
        "求助卡不教“把验证码发给我”",
        "不读短信、通讯录、定位；不做远程控制",
    ], size=14)
    add_textbox(slide, 1.0, 6.7, 11.2, 0.24, "架构原则：风险逻辑放在 domain 层，避免页面组件复制判断导致脆弱和冗余。", size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide, 7)

    # 8 Progress
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_title(slide, "当前进展与下一步", "MVP 已上线，下一阶段需要真实用户验证。")
    add_round_rect(slide, 0.84, 1.6, 5.45, 4.75, WHITE, LINE)
    add_textbox(slide, 1.16, 1.95, 4.7, 0.34, "已完成", size=20, color=NAVY, bold=True)
    add_bullet_list(slide, 1.16, 2.55, 4.7, 2.9, [
        "核心领域模型：risk / question / tutorial / help / routing",
        "首页输入、语音、确认页",
        "低风险分步教程",
        "高风险中断 + 家人求助卡",
        "AI 风险复检与 Vercel 部署",
    ], size=13)
    add_round_rect(slide, 6.85, 1.6, 5.45, 4.75, WHITE, LINE)
    add_textbox(slide, 7.17, 1.95, 4.7, 0.34, "下一步", size=20, color=NAVY, bold=True)
    add_bullet_list(slide, 7.17, 2.55, 4.7, 2.9, [
        "真实老人 / 家属可用性测试",
        "扩充白名单教程库",
        "优化求助卡分享体验",
        "家人端与社区服务站场景",
        "方言 ASR 与截图识别探索",
    ], size=13)
    add_round_rect(slide, 0.84, 6.58, 11.45, 0.52, NAVY, radius=True)
    add_textbox(slide, 1.12, 6.73, 10.9, 0.2, "EasyPhone AI 不只是帮助老人使用智能手机，而是帮助他们更安全地使用智能手机。", size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)

    prs.save(PPTX)
    print(PPTX)


if __name__ == "__main__":
    build_deck()
